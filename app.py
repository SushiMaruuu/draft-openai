import streamlit as st
from pptx import Presentation
import openai
import re
import os
import pymysql
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# Set your OpenAI API key from the environment variable
openai.api_key = os.getenv("OPENAI_API_KEY")

# Function to clean text
def clean_text(text):
    course_number_pattern = r"\b[A-Z]{2,4}\d{4,6}\b"
    lecturer_info_pattern = r"(Prof\.?|Dr\.?|Lecturer|Professor|Mr\.?|Ms\.?)\s[A-Z][a-z]+"
    
    text = re.sub(course_number_pattern, '', text)
    text = re.sub(lecturer_info_pattern, '', text)
    
    unwanted_phrases = ['Slide', 'OCTOBER', 'Short URL']
    for phrase in unwanted_phrases:
        text = text.replace(phrase, '')
    
    text = text.replace('\u2013', '-')  
    return text

# Function to generate MCQs and summary using OpenAI GPT model
# Function to generate MCQs and summary using OpenAI GPT model
def generate_mcqs_and_summary(text):
    truncated_text = text[:4000]
    
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that generates multiple-choice questions with answers and a concise summary from content."},
                {"role": "user", "content": f"Based on the following content, generate exactly 10 distinct multiple-choice questions with four unique options each (one correct answer). Format as follows:\nQuestion: <text>\nOption A: <text>\nOption B: <text>\nOption C: <text>\nOption D: <text>\nAnswer: Option <correct letter>.\n\nAlso, provide a concise summary of the content below.\n\n{truncated_text}"}
            ],
            max_tokens=1500,
            temperature=0.7,
        )
        
        generated_text = response['choices'][0]['message']['content'].strip()
        print("Generated Text:\n", generated_text)  # Debugging: Print the generated text to inspect
        
        if not generated_text:
            st.error("No response received from the OpenAI API.")
            return [], "No summary generated."

        # Extract questions and summary
        mcqs = []
        summary = ""
        
        # Split the generated text into lines for processing
        lines = generated_text.split("\n")
        current_question = ""
        options = {}

        for line in lines:
            line = line.strip()
            if line.startswith("Question:"):
                if current_question and options:
                    mcqs.append({
                        "question": current_question,
                        "options": options,
                        "answer": answer_text
                    })
                current_question = line[len("Question: "):]
                options = {}
                answer_text = ""
            elif line.startswith("Option A:"):
                options['A'] = line[len("Option A: "):]
            elif line.startswith("Option B:"):
                options['B'] = line[len("Option B: "):]
            elif line.startswith("Option C:"):
                options['C'] = line[len("Option C: "):]
            elif line.startswith("Option D:"):
                options['D'] = line[len("Option D: "):]
            elif line.startswith("Answer:"):
                answer_text = line[len("Answer: "):]
        
        # Add the last question if present
        if current_question and options:
            mcqs.append({
                "question": current_question,
                "options": options,
                "answer": answer_text
            })

        # Find the summary by checking for the last line that starts with "Summary:"
        for line in lines:
            if line.lower().startswith("summary:"):
                summary = line[len("Summary: "):].strip()
                break
            elif line:  # If it's not empty, consider it part of the summary
                summary += line + " "

        return mcqs, summary.strip()

    except Exception as e:
        st.error(f"Error generating MCQs and summary: {str(e)}")
        return [], ""

# Function to save MCQs and summary to the database
def save_to_database(subject, topic, mcqs, summary):
    try:
        connection = pymysql.connect(
            host='localhost',
            user='root',
            password='',
            database='mcq_database'
        )
        with connection.cursor() as cursor:
            for mcq in mcqs:
                sql = """
                    INSERT INTO mcq_questions 
                    (subject, topic, question, option_a, option_b, option_c, option_d, answer) 
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
                """
                cursor.execute(sql, (
                    subject,
                    topic,
                    mcq["question"],
                    mcq["options"]["A"],
                    mcq["options"]["B"],
                    mcq["options"]["C"],
                    mcq["options"]["D"],
                    mcq["answer"]
                ))
            connection.commit()
        st.success("Questions and summary saved to the database successfully!")
    except Exception as e:
        st.error(f"Error saving to database: {str(e)}")
    finally:
        connection.close()

# Streamlit app
st.title("PowerPoint Text Extractor and MCQ Generator")

# Input fields for subject and topic
subject = st.text_input("Enter the Subject")
topic = st.text_input("Enter the Topic")

# File uploader for PPTX files only
uploaded_file = st.file_uploader("Upload a PowerPoint (PPTX)", type=["pptx"])

if uploaded_file is not None:
    st.write(f"Filename: {uploaded_file.name}")
    
    presentation = Presentation(uploaded_file)
    all_text = []

    for i, slide in enumerate(presentation.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        
        cleaned_slide_text = clean_text("\n".join(slide_text))
        all_text.append(f"Slide {i + 1}:\n{cleaned_slide_text}\n")

    st.subheader("Extracted Text from PowerPoint")
    for slide_text in all_text:
        st.write(slide_text)

    full_text = " ".join([clean_text(slide) for slide in all_text])
    st.write("Full Extracted Text:\n", full_text)  # Debugging: Check the content being passed to OpenAI

    # Button to generate multiple-choice questions and summary
    if st.button("Generate MCQs and Summary"):
        mcqs, summary = generate_mcqs_and_summary(full_text)
        
        # Store generated data in session state
        st.session_state['mcqs'] = mcqs
        st.session_state['summary'] = summary

        st.subheader("Generated Multiple-Choice Questions with Answers")
        for mcq in mcqs:
            st.write(f"*Question:* {mcq['question']}")
            for opt, text in mcq['options'].items():
                st.write(f"Option {opt}: {text}")
            st.write(f"*Answer:* {mcq['answer']}")
            st.write("---")

        st.subheader("Generated Summary")
        if summary:
            st.write(summary)

    # Button to save data to database if MCQs are generated
    if 'mcqs' in st.session_state and 'summary' in st.session_state:
        if st.button("Save to Database"):
            save_to_database(subject, topic, st.session_state['mcqs'], st.session_state['summary'])